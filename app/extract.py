# app/extract.py
from pathlib import Path
import pdfplumber
from pptx import Presentation
import docx
import os
import shutil
import subprocess
import tempfile
from typing import Optional

# Optional: image OCR dependencies (pytesseract + PIL)
try:
    from PIL import Image
    import pytesseract

    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False


def extract_text_from_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text_from_pdf(path: Path) -> str:
    text_pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            txt = page.extract_text()
            if txt:
                text_pages.append(f"[page:{i}]\n" + txt)
    return "\n\n".join(text_pages)


def _presentation_to_text(prs: Presentation) -> str:
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            # some shapes may not have .text
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append(f"[slide:{i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _convert_ppt_to_pptx_with_libreoffice(src_path: Path, timeout: int = 120) -> Path:
    """
    Convert a legacy .ppt (binary) file to .pptx using LibreOffice (soffice).
    Returns the Path to the converted .pptx file.

    The converted file is produced in a temporary directory and returned to the caller.
    """
    soffice = shutil.which("soffice")
    if soffice is None:
        raise RuntimeError(
            "LibreOffice 'soffice' not found in PATH. Install LibreOffice (brew install --cask libreoffice) "
            "or convert the .ppt to .pptx manually."
        )

    tmpdir = Path(tempfile.mkdtemp(prefix="ppt_convert_"))
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pptx",
        "--outdir",
        str(tmpdir),
        str(src_path),
    ]

    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
    if proc.returncode != 0:
        stderr = proc.stderr.decode("utf-8", errors="ignore").strip()
        raise RuntimeError(f"LibreOffice conversion failed: {stderr[:1000]}")

    converted = tmpdir / (src_path.stem + ".pptx")
    if not converted.exists():
        raise RuntimeError("LibreOffice conversion succeeded but output .pptx file not found.")

    return converted


def extract_text_from_pptx(path: Path) -> str:
    """
    Robust extraction for PowerPoint files.

    - If the path points to a .ppt (legacy binary), attempt to convert to .pptx
      using LibreOffice and then extract text.
    - If the path points to a .pptx, parse directly with python-pptx.

    Raises RuntimeError with actionable messages on failure.
    """
    suf = path.suffix.lower()
    if suf == ".ppt":
        # convert -> extract
        try:
            converted = _convert_ppt_to_pptx_with_libreoffice(path)
            try:
                prs = Presentation(str(converted))
                return _presentation_to_text(prs)
            except Exception as e:
                raise RuntimeError(f"Extraction from converted PPTX failed: {e}")
        except Exception as e:
            raise RuntimeError(f"PPT -> PPTX conversion/extraction failed: {e}")

    # assume .pptx
    try:
        prs = Presentation(str(path))
        return _presentation_to_text(prs)
    except Exception as e:
        raise RuntimeError(
            f"Failed reading PPTX: {e}. If the file is a legacy .ppt, install LibreOffice and enable conversion, "
            "or open & re-save as .pptx in PowerPoint and re-upload."
        )


def extract_text_from_docx(path: Path) -> str:
    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_image(path: Path) -> str:
    if not OCR_AVAILABLE:
        raise RuntimeError(
            "OCR not available. Install pytesseract and pillow, and ensure tesseract is installed on system."
        )
    img = Image.open(path)
    return pytesseract.image_to_string(img)


def extract_text(filepath: str) -> str:
    """
    Generic extractor: chooses method by extension.
    Returns extracted text string (may include page/slide markers).
    """
    p = Path(filepath)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    ext = p.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(p)
    if ext == ".pptx":
        return extract_text_from_pptx(p)
    if ext == ".ppt":
        # convert .ppt -> .pptx then extract
        converted = _convert_ppt_to_pptx_with_libreoffice(p)
        return extract_text_from_pptx(converted)
    if ext in {".docx", ".doc"}:
        return extract_text_from_docx(p)
    if ext in {".txt", ".md"}:
        return extract_text_from_txt(p)
    if ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}:
        return extract_text_from_image(p)

    raise ValueError(f"Unsupported file extension: {ext}")